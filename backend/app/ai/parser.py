"""多格式文件解析器。

将 PPT/PDF/Word/Excel/TXT/MD 统一解析为纯文本，供 AI 训练使用。
管理员上传任意格式文件，系统自动解析转换，无需关心中间过程。

支持的格式：
- TXT / MD / JSON: 直接读取文本
- DOCX: 提取段落 + 表格
- PPTX: 提取每页标题 + 正文 + 备注
- PDF: 提取文本（图片 OCR 第二阶段）
- XLSX: 提取单元格数据（按行输出）
"""
import logging
from pathlib import Path
from typing import Callable

logger = logging.getLogger(__name__)

# 支持的文件扩展名
SUPPORTED_EXTENSIONS = {"txt", "md", "json", "docx", "pptx", "pdf", "xlsx"}


def parse_file(file_path: Path, file_type: str = None) -> str:
    """根据文件类型解析为纯文本。

    Args:
        file_path: 文件路径
        file_type: 文件扩展名（不含.），如不传则从路径推断

    Returns:
        解析后的纯文本
    """
    if file_type is None:
        file_type = file_path.suffix.lstrip(".").lower()

    parsers: dict[str, Callable] = {
        "txt": _parse_text,
        "md": _parse_text,
        "json": _parse_text,
        "docx": _parse_docx,
        "pptx": _parse_pptx,
        "pdf": _parse_pdf,
        "xlsx": _parse_xlsx,
    }

    parser = parsers.get(file_type)
    if parser is None:
        logger.warning("不支持的文件类型: %s，按文本处理", file_type)
        return _parse_text(file_path)

    try:
        text = parser(file_path)
        logger.info("解析 %s 完成，文本长度 %d 字符", file_path.name, len(text))
        return text
    except Exception as e:
        logger.error("解析 %s 失败: %s", file_path.name, e)
        raise ValueError(f"文件解析失败: {e}")


def _parse_text(file_path: Path) -> str:
    """TXT/MD/JSON 直接读取。"""
    # 尝试 utf-8，失败用 gbk
    try:
        return file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return file_path.read_text(encoding="gbk", errors="ignore")


def _parse_docx(file_path: Path) -> str:
    """Word 文档：提取段落 + 表格。"""
    from docx import Document

    doc = Document(str(file_path))
    parts = []

    # 段落
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 表格
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts) if parts else ""


def _parse_pptx(file_path: Path) -> str:
    """PPT：提取每页标题 + 正文 + 备注。"""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []

        # 标题
        if slide.shapes.title and slide.shapes.title.text:
            slide_parts.append(f"【第{i}页】{slide.shapes.title.text.strip()}")

        # 正文（文本框）
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text not in (slide_parts[0] if slide_parts else ""):
                    slide_parts.append(text)

        # 备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"备注：{notes}")

        if slide_parts:
            parts.append("\n".join(slide_parts))

    return "\n\n---\n\n".join(parts) if parts else ""


def _parse_pdf(file_path: Path) -> str:
    """PDF：提取文本（图片 OCR 第二阶段）。"""
    import pdfplumber

    parts = []
    with pdfplumber.open(str(file_path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                parts.append(f"【第{i}页】\n{text.strip()}")

    return "\n\n---\n\n".join(parts) if parts else ""


def _parse_xlsx(file_path: Path) -> str:
    """Excel：按 sheet 和行提取数据。"""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"【工作表：{sheet_name}】\n" + "\n".join(rows))

    wb.close()
    return "\n\n---\n\n".join(parts) if parts else ""
